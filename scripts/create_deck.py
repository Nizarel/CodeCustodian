"""Generate CodeCustodian 2-slide PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import os

# ── Generate Copilot SDK banner image ──
def create_copilot_sdk_banner(path: str, width: int = 800, height: int = 200):
    """Create a GitHub Copilot SDK branded banner."""
    img = Image.new("RGBA", (width, height), (13, 17, 23, 255))
    draw = ImageDraw.Draw(img)

    # Rounded rect border
    border_color = (88, 166, 255, 180)
    draw.rounded_rectangle(
        [4, 4, width - 5, height - 5], radius=16, outline=border_color, width=2
    )

    # Try to use a nice font, fall back to default
    try:
        font_large = ImageFont.truetype("segoeui.ttf", 52)
        font_medium = ImageFont.truetype("segoeui.ttf", 30)
        font_small = ImageFont.truetype("segoeuib.ttf", 18)
    except OSError:
        try:
            font_large = ImageFont.truetype("arial.ttf", 52)
            font_medium = ImageFont.truetype("arial.ttf", 30)
            font_small = ImageFont.truetype("arialbd.ttf", 18)
        except OSError:
            font_large = ImageFont.load_default()
            font_medium = font_large
            font_small = font_large

    # "GitHub" text
    draw.text((40, 25), "GitHub", fill=(63, 185, 80), font=font_small)

    # "Copilot" in purple gradient simulation
    draw.text((40, 50), "COPILOT", fill=(188, 140, 255), font=font_medium)

    # "SDK" in green gradient simulation
    draw.text((40, 90), "SDK", fill=(63, 185, 80), font=font_large)

    # Robot face (pixel art style) on the right
    rx, ry = width - 180, 20
    size = 12

    # Head outline (blue)
    head_color = (88, 166, 255, 255)
    purple = (188, 140, 255, 255)
    dark = (13, 17, 23, 255)

    # Antenna
    draw.rectangle([rx + 60, ry, rx + 72, ry + 12], fill=purple)
    draw.rectangle([rx + 72, ry, rx + 84, ry + 12], fill=(200, 130, 220))

    # Head top row
    for i in range(10):
        draw.rectangle(
            [rx + 24 + i * size, ry + 12, rx + 36 + i * size, ry + 24],
            fill=head_color,
        )

    # Head body (8 rows)
    for row in range(8):
        y = ry + 24 + row * size
        # Left edge
        draw.rectangle([rx + 12, y, rx + 24, y + size], fill=head_color)
        # Right edge
        draw.rectangle([rx + 144, y, rx + 156, y + size], fill=head_color)
        # Fill
        for col in range(10):
            x = rx + 24 + col * size
            draw.rectangle([x, y, x + size, y + size], fill=(40, 60, 100))

    # Eyes (big white squares with dark pupils)
    for ex in [rx + 36, rx + 96]:
        # White of eye
        for er in range(3):
            for ec in range(3):
                draw.rectangle(
                    [ex + ec * size, ry + 48 + er * size,
                     ex + (ec + 1) * size, ry + 48 + (er + 1) * size],
                    fill=(200, 220, 255),
                )
        # Pupil
        draw.rectangle([ex + size, ry + 48 + size, ex + 2 * size, ry + 48 + 2 * size],
                        fill=dark)

    # Mouth
    for i in range(4):
        draw.rectangle(
            [rx + 48 + i * size, ry + 108, rx + 60 + i * size, ry + 120],
            fill=head_color,
        )

    # Ears (purple)
    for ey_off in range(3):
        # Left ear
        draw.rectangle(
            [rx, ry + 48 + ey_off * size, rx + 12, ry + 60 + ey_off * size],
            fill=purple,
        )
        # Right ear
        draw.rectangle(
            [rx + 156, ry + 48 + ey_off * size, rx + 168, ry + 60 + ey_off * size],
            fill=purple,
        )

    # Bottom row
    for i in range(10):
        draw.rectangle(
            [rx + 24 + i * size, ry + 132, rx + 36 + i * size, ry + 144],
            fill=head_color,
        )

    img.save(path, "PNG")
    return path


BANNER_PATH = "presentations/copilot-sdk-banner.png"
create_copilot_sdk_banner(BANNER_PATH)
print(f"✅ Generated {BANNER_PATH}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG = RGBColor(0x0D, 0x11, 0x17)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
BLUE = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
ORANGE = RGBColor(0xD2, 0x99, 0x22)
RED = RGBColor(0xF8, 0x51, 0x49)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
DARK_BG = RGBColor(0x16, 0x1B, 0x22)
CARD_BG = RGBColor(0x21, 0x26, 0x2D)
BORDER = RGBColor(0x30, 0x36, 0x3D)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_metric_card(slide, left, top, number, label, num_color):
    w, h = Inches(2.6), Inches(2.2)
    add_rounded_rect(slide, left, top, w, h, CARD_BG, BORDER)
    add_text_box(slide, left, top + Inches(0.3), w, Inches(1),
                 number, font_size=40, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.2), w - Inches(0.3), Inches(0.9),
                 label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


def add_pipeline_step(slide, left, top, text, border_color=None):
    w, h = Inches(1.35), Inches(0.45)
    bc = border_color or BORDER
    add_rounded_rect(slide, left, top, w, h, CARD_BG, bc)
    add_text_box(slide, left, top + Inches(0.05), w, h,
                 text, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)


def add_arrow(slide, left, top):
    add_text_box(slide, left, top + Inches(0.05), Inches(0.3), Inches(0.4),
                 "→", font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Business Value Proposition
# ═══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG)

# Title
add_text_box(slide1, Inches(0), Inches(0.3), prs.slide_width, Inches(0.9),
             "CodeCustodian", font_size=48, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(0), Inches(1.1), prs.slide_width, Inches(0.5),
             "Autonomous AI Agent for Technical Debt Management", font_size=20,
             color=GRAY, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide1, Inches(0), Inches(1.6), prs.slide_width, Inches(0.5),
             '"Engineering teams spend 40% on maintenance. What if an AI agent handled it?"',
             font_size=14, color=BLUE, alignment=PP_ALIGN.CENTER)

# Badge row
badges = ["Python 3.11+", "GitHub Copilot SDK", "FastMCP v2", "Azure Container Apps", "953 Tests · 82% Cov"]
badge_colors = [GREEN, BLUE, PURPLE, AZURE_BLUE, GREEN]
bx = Inches(1.8)
for i, (badge, bc) in enumerate(zip(badges, badge_colors)):
    add_rounded_rect(slide1, bx, Inches(2.15), Inches(1.85), Inches(0.35), CARD_BG, bc)
    add_text_box(slide1, bx, Inches(2.18), Inches(1.85), Inches(0.35),
                 badge, font_size=10, color=bc, alignment=PP_ALIGN.CENTER)
    bx += Inches(1.95)

# Metric cards
card_x = Inches(0.65)
card_y = Inches(2.85)
add_metric_card(slide1, card_x, card_y, "7", "Built-in Scanners\n(Security, APIs, TODOs,\nSmells, Types, Drift, Deps)", GREEN)
add_metric_card(slide1, card_x + Inches(2.85), card_y, "$4,960", "Customer-validated\nsavings (3-week pilot)", BLUE)
add_metric_card(slide1, card_x + Inches(5.7), card_y, "62h", "Engineering hours saved\n(per 3-week cycle)", PURPLE)
add_metric_card(slide1, card_x + Inches(8.55), card_y, "17", "MCP Tools + 7 Prompts\n+ 12 Agent Profiles", ORANGE)

# Pipeline
py = Inches(5.3)
px = Inches(0.7)
steps = ["Scan", "De-dup", "Prioritize", "Plan (SDK)", "Execute", "Verify", "PR"]
for i, step in enumerate(steps):
    add_pipeline_step(slide1, px, py, step)
    if i < len(steps) - 1:
        add_arrow(slide1, px + Inches(1.35), py)
    px += Inches(1.65)

# Confidence gates
gate_y = Inches(5.95)
gates = [("Confidence 8-10 → Auto-PR", GREEN), ("Confidence 5-7 → Draft PR", ORANGE), ("Confidence <5 → Proposal Only", RED)]
gx = Inches(2.5)
for text, color in gates:
    add_rounded_rect(slide1, gx, gate_y, Inches(2.7), Inches(0.4), CARD_BG, color)
    add_text_box(slide1, gx, gate_y + Inches(0.04), Inches(2.7), Inches(0.35),
                 text, font_size=10, color=color, alignment=PP_ALIGN.CENTER)
    gx += Inches(2.85)

# Repo link
add_text_box(slide1, Inches(0), Inches(6.8), prs.slide_width, Inches(0.4),
             "https://github.com/Nizarel/CodeCustodian", font_size=12,
             color=BLUE, alignment=PP_ALIGN.CENTER)

# Copilot SDK banner (bottom-right corner of slide 1)
if os.path.exists(BANNER_PATH):
    slide1.shapes.add_picture(BANNER_PATH, Inches(9.3), Inches(6.55), Inches(3.6), Inches(0.9))

add_text_box(slide1, Inches(0), Inches(7.1), prs.slide_width, Inches(0.3),
             "v0.15.2 · MIT License · Python 3.11+", font_size=10,
             color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture & Azure Integration
# ═══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide2, BG)

# Title
add_text_box(slide2, Inches(0), Inches(0.2), prs.slide_width, Inches(0.7),
             "Architecture & Azure Integration", font_size=36, color=BLUE,
             bold=True, alignment=PP_ALIGN.CENTER)


def add_arch_box(slide, left, top, width, height, title, items, title_color, border_color):
    add_rounded_rect(slide, left, top, width, height, DARK_BG, border_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
                 title, font_size=15, color=title_color, bold=True)
    y_off = top + Inches(0.5)
    for item in items:
        add_text_box(slide, left + Inches(0.35), y_off, width - Inches(0.6), Inches(0.28),
                     f"→ {item}", font_size=10, color=GRAY)
        y_off += Inches(0.28)


box_w = Inches(5.9)
box_h_top = Inches(2.5)
box_h_bot = Inches(2.2)
col1 = Inches(0.5)
col2 = Inches(6.9)
row1 = Inches(1.05)
row2 = Inches(3.75)

# AI & Dev Experience (top-left)
add_arch_box(slide2, col1, row1, box_w, box_h_top,
             "🤖 AI & Developer Experience", [
                 "GitHub Copilot SDK — AI-powered refactoring (12 agent profiles)",
                 "FastMCP v2 Server — 17 tools, 7 prompts, 8 resources",
                 "Typer + Rich CLI — 15 commands (JSON/table/CSV/SARIF)",
                 "AI Test Synthesis — auto-generate regression tests",
                 "Agentic Migrations — multi-stage framework upgrades",
                 "Remote Repo Scanning — scan any public Git repo via URL",
                 "Feedback Loop — learns from PR outcomes over time",
             ], PURPLE, PURPLE)

# Azure (top-right)
add_arch_box(slide2, col2, row1, box_w, box_h_top,
             "☁️ Deployed on Azure (7 services)", [
                 "Azure Key Vault — secrets management (GitHub PAT, webhooks)",
                 "Azure Container Apps — serverless MCP deployment (Bicep IaC)",
                 "Azure Monitor — OpenTelemetry observability",
                 "Azure DevOps — work item integration",
                 "Microsoft Teams — ChatOps Adaptive Cards + Work IQ",
                 "GitHub Actions — 6 CI/CD workflows (lint, test, deploy, scan)",
                 "Work IQ MCP — sprint context enrichment",
             ], AZURE_BLUE, AZURE_BLUE)

# Enterprise (bottom-left)
add_arch_box(slide2, col1, row2, box_w, box_h_bot,
             "🏢 Enterprise Features", [
                 "Budget Manager — prevents AI spend overruns",
                 "SLA Reporter — debt resolution tracking",
                 "ROI Calculator — cost savings with HTML/CSV export",
                 "RBAC + Approval Workflows — multi-team governance",
                 "Business Impact Scoring — 5-factor prioritization",
                 "Predictive Debt Forecasting — trend analysis",
             ], ORANGE, ORANGE)

# Security (bottom-right)
add_arch_box(slide2, col2, row2, box_w, box_h_bot,
             "🛡️ Security & Responsible AI", [
                 "7-point safety system (syntax, size, path, encoding, secrets, blast radius)",
                 "Audit Logger — SHA-256 tamper-evident chain",
                 "Dangerous function detection (eval, exec, __import__)",
                 "Confidence-gated safety — low confidence = proposal only",
                 "Responsible AI Policy — 8 documented principles",
                 "Bandit + Trivy security scanning in CI/CD",
             ], RED, RED)

# Footer stats
footer_y = Inches(6.2)
stats = [("953", "tests"), ("82%", "coverage"), ("6", "CI/CD workflows"), ("0", "lint errors"), ("v0.15.2", "")]
fx = Inches(2.2)
for val, label in stats:
    text = f"{val} {label}" if label else val
    add_text_box(slide2, fx, footer_y, Inches(1.6), Inches(0.3),
                 text, font_size=11, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    fx += Inches(1.8)

# Repo link + pipeline flow
add_text_box(slide2, Inches(0), Inches(6.65), prs.slide_width, Inches(0.35),
             "Scan → De-dup → Prioritize → Plan (Copilot SDK) → Execute → Verify → PR",
             font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide2, Inches(0), Inches(7.05), prs.slide_width, Inches(0.35),
             "https://github.com/Nizarel/CodeCustodian", font_size=12,
             color=BLUE, alignment=PP_ALIGN.CENTER)

# Copilot SDK banner (bottom-right corner of slide 2)
if os.path.exists(BANNER_PATH):
    slide2.shapes.add_picture(BANNER_PATH, Inches(9.3), Inches(6.55), Inches(3.6), Inches(0.9))


# ── Save ──
output = "presentations/CodeCustodian-Deck-v2.pptx"
prs.save(output)
print(f"✅ Saved {output}")
