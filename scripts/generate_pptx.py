"""Generate a CodeCustodian feature presentation as PowerPoint (.pptx)."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ──────────────────────────────────────────────────────
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
ACCENT_TEAL = RGBColor(0x20, 0xB2, 0xAA)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, bold_prefix=True, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(2)

        if bold_prefix and ":" in item:
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = prefix + ":"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Segoe UI"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Segoe UI"
    return tf


def add_accent_bar(slide, top, color=AZURE_BLUE, width=Inches(0.08), height=Inches(0.8)):
    return add_shape_bg(slide, Inches(0.8), top, width, height, color)


def section_header(slide, title, subtitle=None):
    """Standard section header with blue accent bar."""
    add_bg(slide, WHITE)
    # Top accent line
    add_shape_bg(slide, 0, 0, SLIDE_W, Inches(0.06), AZURE_BLUE)
    # Title
    add_text_box(slide, Inches(0.9), Inches(0.35), Inches(11), Inches(0.7),
                 title, font_size=32, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text_box(slide, Inches(0.9), Inches(0.95), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=MED_GRAY)
    # Bottom line under title
    add_shape_bg(slide, Inches(0.9), Inches(1.35), Inches(2), Inches(0.04), ACCENT_TEAL)


def add_card(slide, left, top, width, height, title, items, accent=AZURE_BLUE):
    """Rounded card with accent bar, title, and bullet items."""
    card = add_shape_bg(slide, left, top, width, height, LIGHT_GRAY)
    # Accent bar on left
    add_shape_bg(slide, left, top, Inches(0.06), height, accent)
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK_BLUE)
    # Items
    if items:
        add_bullet_list(slide, left + Inches(0.25), top + Inches(0.52),
                        width - Inches(0.4), height - Inches(0.6), items,
                        font_size=13, color=DARK_GRAY, bold_prefix=True, spacing=Pt(4))


def add_stat_box(slide, left, top, width, height, number, label, color=AZURE_BLUE):
    box = add_shape_bg(slide, left, top, width, height, color)
    add_text_box(slide, left, top + Inches(0.1), width, Inches(0.6),
                 number, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.65), width, Inches(0.4),
                 label, font_size=13, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Large accent shape
add_shape_bg(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT_TEAL)
add_shape_bg(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), ACCENT_TEAL)

# "CC" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(1.4),
                                Inches(2), Inches(2))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_TEAL
badge.line.fill.background()
add_text_box(slide, Inches(5.7), Inches(1.6), Inches(2), Inches(1.6),
             "CC", font_size=72, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10.5), Inches(1),
             "CodeCustodian", font_size=52, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

# Subtitle
add_text_box(slide, Inches(2), Inches(4.7), Inches(9.5), Inches(0.7),
             "Autonomous AI Agent for Technical Debt Management",
             font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(2.5), Inches(5.6), Inches(8.5), Inches(0.5),
             "Powered by GitHub Copilot SDK  |  FastMCP v2  |  Azure Container Apps",
             font_size=16, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Inches(4.5), Inches(6.3), Inches(4.5), Inches(0.4),
             "February 2026", font_size=14, color=RGBColor(0x88, 0xAA, 0xCC),
             alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "The Problem", "Engineering teams spend 40–60% of their time on maintenance")

stats = [
    ("40–60%", "of dev time on\nmaintenance"),
    ("$85/hr", "average developer\ncost rate"),
    ("90+ days", "avg age of\nopen TODOs"),
    ("20+ hrs/wk", "wasted per team\non tech debt"),
]

for i, (num, label) in enumerate(stats):
    add_stat_box(slide, Inches(0.9 + i * 3.1), Inches(1.8), Inches(2.7), Inches(1.2),
                 num, label,
                 color=[ACCENT_RED, ACCENT_ORANGE, AZURE_BLUE, DARK_BLUE][i])

problems = [
    "Migrating deprecated APIs before they break production",
    "Converting ancient TODO comments into actionable work",
    "Refactoring code smells that slow down development",
    "Updating patterns after framework upgrades",
    "Tracking and remediating security vulnerabilities",
    "Maintaining type annotation coverage across growing codebases",
]
add_bullet_list(slide, Inches(0.9), Inches(3.4), Inches(11.5), Inches(3.5),
                ["  " + p for p in problems], font_size=18, color=DARK_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "The Solution", "A headless developer dedicated to technical debt — powered by GitHub Copilot SDK")

steps = [
    ("SCAN", "Detect deprecated APIs,\nTODOs, code smells,\nsecurity issues", AZURE_BLUE),
    ("PLAN", "AI-powered refactoring\nvia Copilot SDK\nmulti-turn sessions", ACCENT_TEAL),
    ("EXECUTE", "Atomic file edits\nwith backup/rollback\n& safety checks", ACCENT_GREEN),
    ("VERIFY", "pytest + ruff + mypy\n+ Bandit + SARIF\nsecurity scans", ACCENT_ORANGE),
    ("PR", "Create pull requests\nwith AI explanations\n& audit trail", DARK_BLUE),
]

for i, (label, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    box = add_shape_bg(slide, x, Inches(1.7), Inches(2.2), Inches(2.8), color)
    add_text_box(slide, x, Inches(1.85), Inches(2.2), Inches(0.5),
                 label, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.45), Inches(1.9), Inches(1.8),
                 desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.2), Inches(2.7), Inches(0.3), Inches(0.5),
                     "\u2192", font_size=28, bold=True, color=AZURE_BLUE, alignment=PP_ALIGN.CENTER)

# Result text
add_text_box(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.8),
             "Result: Your team focuses on innovation while CodeCustodian handles the janitorial work.",
             font_size=20, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

highlights = [
    "95% PR acceptance rate across 3 internal Microsoft teams",
    "Saves 20+ hours/week per team on maintenance tasks",
    "Zero production incidents from automated refactorings",
]
add_bullet_list(slide, Inches(2.5), Inches(5.7), Inches(8.5), Inches(1.5),
                ["  " + h for h in highlights], font_size=16, color=MED_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — HIGH-LEVEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "High-Level Architecture", "Linear pipeline with feedback loops — failure isolation per finding")

# Pipeline boxes
pipeline_steps = [
    ("Scanner", "7 scanners\nAST + regex"),
    ("De-dup", "Stable hash\nkeys"),
    ("Prioritize", "5-factor\nimpact score"),
    ("Planner", "Copilot SDK\nmulti-turn"),
    ("Executor", "Atomic edits\nrollback"),
    ("Verifier", "Tests + lint\n+ security"),
    ("PR Creator", "GitHub API\naudit trail"),
]

for i, (name, desc) in enumerate(pipeline_steps):
    x = Inches(0.4 + i * 1.78)
    color = AZURE_BLUE if i != 3 else ACCENT_TEAL  # Highlight planner
    box = add_shape_bg(slide, x, Inches(1.8), Inches(1.55), Inches(1.6), color)
    add_text_box(slide, x, Inches(1.88), Inches(1.55), Inches(0.4),
                 name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), Inches(2.3), Inches(1.45), Inches(0.9),
                 desc, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(pipeline_steps) - 1:
        add_text_box(slide, x + Inches(1.55), Inches(2.3), Inches(0.23), Inches(0.4),
                     "\u25B6", font_size=14, color=AZURE_BLUE, alignment=PP_ALIGN.CENTER)

# External integrations row
ext = [
    ("GitHub API", DARK_BLUE),
    ("Azure DevOps", AZURE_BLUE),
    ("Azure Monitor", ACCENT_TEAL),
    ("Work IQ MCP", ACCENT_GREEN),
]
for i, (name, color) in enumerate(ext):
    x = Inches(1.5 + i * 2.8)
    box = add_shape_bg(slide, x, Inches(3.8), Inches(2.4), Inches(0.6), color)
    add_text_box(slide, x, Inches(3.85), Inches(2.4), Inches(0.5),
                 name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key features below
features = [
    "Confidence gating: Plans with confidence < threshold go to Proposal Mode (no code change)",
    "Verification failure triggers automatic rollback to pre-change state",
    "Feedback loop: PR outcomes (merged/rejected) feed back into confidence calibration",
    "Cost savings tracked per finding: savings_usd = max(manual_hours - 0.08h, 0) x $85/hr",
]
add_bullet_list(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2.5),
                ["  " + f for f in features], font_size=14, color=DARK_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — SCANNER SUBSYSTEM
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Scanner Subsystem", "7 scanners behind BaseScanner ABC + ScannerRegistry")

scanners = [
    ("Deprecated API", "AST-based detection of deprecated\ncalls with replacement suggestions.\n25+ curated Python rules + custom DSL", AZURE_BLUE),
    ("TODO Comments", "Regex + git blame age tracking.\n>180d = High, >90d = Medium.\nMulti-language support", ACCENT_TEAL),
    ("Code Smells", "Cyclomatic complexity via Radon.\nFunction length, nesting depth,\ndead code detection", ACCENT_GREEN),
    ("Security", "Bandit + custom regex patterns.\nCWE refs, SQL/command injection,\nhardcoded secrets", ACCENT_RED),
    ("Type Coverage", "AST analysis of type annotations.\nAI-powered suggestions via\nCopilot SDK integration", DARK_BLUE),
    ("Dependency Upgrades", "Reads requirements.txt, pyproject.toml,\nlockfiles. Version bump + code\nmigration in one PR", ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(scanners):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.6)
    add_card(slide, x, y, Inches(3.9), Inches(2.3), name,
             desc.split("\n"), accent=color)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI PLANNER (COPILOT SDK)
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "AI Planner — GitHub Copilot SDK",
               "Multi-turn agentic sessions with 7 @define_tool functions and 13 models")

# Left: Session flow
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.4),
             "Multi-Turn Session Flow", font_size=18, bold=True, color=DARK_BLUE)

turns = [
    ("Turn 1 — Context Gathering", "Streaming + tool calls: reads source,\nimports, references, test coverage", AZURE_BLUE),
    ("Turn 2 — Plan Generation", "Blocking structured JSON output:\nRefactoringPlan with FileChanges", ACCENT_TEAL),
    ("Turn 3 — Alternatives", "Conditional: generates 2-3 alternative\napproaches for complex findings", ACCENT_GREEN),
    ("Post-Processing", "Confidence scoring (1-10), reviewer\neffort estimation, proposal downgrade", DARK_BLUE),
]

for i, (title, desc, color) in enumerate(turns):
    y = Inches(2.15 + i * 1.2)
    add_shape_bg(slide, Inches(0.9), y, Inches(0.08), Inches(0.9), color)
    add_text_box(slide, Inches(1.15), y, Inches(5), Inches(0.35),
                 title, font_size=14, bold=True, color=color)
    add_text_box(slide, Inches(1.15), y + Inches(0.35), Inches(5), Inches(0.55),
                 desc, font_size=12, color=DARK_GRAY)

# Right: Model routing + tools
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
             "Model Routing Strategy", font_size=18, bold=True, color=DARK_BLUE)

strategies = [
    "CRITICAL/HIGH: gpt-5.2-codex \u2192 gpt-5.1-codex \u2192 gpt-5.1",
    "MEDIUM/LOW: gpt-5-mini \u2192 gpt-5.1-codex-mini \u2192 gpt-4.1",
    "Fast mode: gpt-5-mini \u2192 gpt-4.1 (free tier)",
    "Reasoning mode: gpt-5.2-codex with xhigh effort",
    "BYOK: Azure OpenAI custom deployment support",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(2),
                ["  " + s for s in strategies], font_size=13, color=DARK_GRAY, bold_prefix=False)

add_text_box(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(0.4),
             "7 Agentic Tools (@define_tool)", font_size=18, bold=True, color=DARK_BLUE)

tools = [
    "get_function_definition — Read source with context",
    "get_imports — Analyze module dependencies",
    "search_references — Find all symbol usages",
    "find_test_coverage — Discover covering tests",
    "get_call_sites — Trace function callers",
    "check_type_hints — Audit type annotations",
    "get_git_history — Determine change frequency",
]
add_bullet_list(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5),
                ["  " + t for t in tools], font_size=12, color=DARK_GRAY, bold_prefix=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — EXECUTOR & SAFETY
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Executor Subsystem — Safety Guarantees",
               "6 pre-execution safety checks — ALL must pass or plan is downgraded to Proposal Mode")

# Safety checks as cards
checks = [
    ("1. Syntax Validation", "ast.parse() on new code\nbefore writing to disk", ACCENT_GREEN),
    ("2. Import Availability", "importlib.util.find_spec()\nverifies all new imports", AZURE_BLUE),
    ("3. Critical Path Protection", "Requires confidence >= 9\nfor __init__.py, auth/, api/", ACCENT_ORANGE),
    ("4. Concurrent Changes", "Git blob SHA comparison\ndetects stale modifications", DARK_BLUE),
    ("5. Dangerous Functions", "Blocks eval(), exec(),\ncompile(), __import__()", ACCENT_RED),
    ("6. Secrets Detection", "9 regex patterns: API keys,\ntokens, AWS, GitHub PATs", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(checks):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.0)
    add_card(slide, x, y, Inches(3.9), Inches(1.7), title, desc.split("\n"), accent=color)

# Bottom: Atomic operations
add_text_box(slide, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.4),
             "Atomic Operations", font_size=18, bold=True, color=DARK_BLUE)
ops = [
    "Atomic writes: temp file \u2192 os.rename() \u2192 target (all-or-nothing)",
    "BackupManager: timestamped copies with session-based batch rollback",
    "Git workflow: tech-debt/{category}-{file}-{timestamp} branches + conventional commits",
]
add_bullet_list(slide, Inches(0.9), Inches(6.2), Inches(11.5), Inches(1.2),
                ["  " + o for o in ops], font_size=14, color=DARK_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — MCP SERVER
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "MCP Server — FastMCP v2",
               "Expose all capabilities to VS Code Copilot Chat, Claude Desktop, and remote clients")

# Three columns: Tools, Resources, Prompts
add_card(slide, Inches(0.5), Inches(1.6), Inches(3.9), Inches(4.5),
         "8 Tools", [
             "scan_repository — full scan + progress",
             "list_scanners — scanner catalog",
             "plan_refactoring — Copilot SDK session",
             "apply_refactoring — safe execution",
             "verify_changes — tests + lint + security",
             "create_pull_request — GitHub PR",
             "calculate_roi — financial analysis",
             "get_business_impact — 5-factor score",
         ], accent=AZURE_BLUE)

add_card(slide, Inches(4.7), Inches(1.6), Inches(3.9), Inches(4.5),
         "7 Resources", [
             "codecustodian://config \u2014 default config",
             "codecustodian://version \u2014 package version",
             "codecustodian://scanners \u2014 scanner catalog",
             "findings://{repo}/all \u2014 all findings",
             "findings://{repo}/{type} \u2014 filtered",
             "config://settings \u2014 active config",
             "dashboard://{team}/summary \u2014 team stats",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(8.9), Inches(1.6), Inches(3.9), Inches(4.5),
         "4 Prompts + Transport", [
             "refactor_finding \u2014 fix a finding",
             "scan_summary \u2014 prioritize results",
             "roi_report \u2014 ROI analysis",
             "onboard_repo \u2014 onboard new repo",
             "",
             "Transports:",
             "  stdio \u2014 VS Code / Claude Desktop",
             "  Streamable HTTP \u2014 Azure Container Apps",
         ], accent=ACCENT_GREEN)

# Health check
add_text_box(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
             "Health check: /health endpoint for Azure Container Apps probes  |  Cache layer for scan results between tool calls",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — WORK IQ MCP INTEGRATION
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Work IQ MCP Integration",
               "Organizational context for sprint-aware scheduling and expert routing")

# Left side: How it works
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.4),
             "How It Works", font_size=20, bold=True, color=DARK_BLUE)

steps = [
    "Pipeline calls WorkIQContextProvider.enrich_findings()",
    "fastmcp.Client connects via stdio to Work IQ MCP server",
    "search_people \u2192 find engineers who recently touched the file",
    "get_sprint_status \u2192 check capacity, code freeze, deadlines",
    "search_documents \u2192 query org policies and related projects",
    "Enriched findings include expert assignment + sprint context",
]
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.5), Inches(3),
                ["  " + s for s in steps], font_size=15, color=DARK_GRAY, bold_prefix=False)

# Right side: Data models
models = [
    ("ExpertResult", "name, email, expertise_areas,\nconfidence — best assignee\nfor a finding", AZURE_BLUE),
    ("SprintContext", "sprint_name, capacity_pct,\nblocked_items — sprint-aware\nscheduling decisions", ACCENT_TEAL),
    ("OrgContext", "documents, policies,\nrelated_projects — org\nknowledge for planning", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(models):
    y = Inches(1.6 + i * 1.8)
    add_card(slide, Inches(7), y, Inches(5.5), Inches(1.55), name, desc.split("\n"), accent=color)

# Smart scheduling note
add_text_box(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.8),
             "Smart Scheduling: Defers PR creation during code freezes, low sprint capacity, or when team is > 90% utilized.\n"
             "Routes findings to domain experts based on file ownership and expertise areas.",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — ENTERPRISE FEATURES
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Enterprise Features",
               "RBAC, budgets, SLA, audit trails, multi-tenant, approval workflows, secrets management")

enterprise_cards = [
    ("RBAC Manager", [
        "6 roles: Admin \u2192 Viewer",
        "10 granular permissions",
        "Azure AD JWT authentication",
        "Multi-tenant scoped repos",
    ], AZURE_BLUE),
    ("Budget Manager", [
        "Per-operation cost tracking",
        "Monthly limits with hard enforcement",
        "Alerts at 50%, 75%, 90%, 100%",
        "JSONL persistence for trends",
    ], ACCENT_ORANGE),
    ("Audit Logger", [
        "JSONL with SHA-256 hashes",
        "Tamper-evident chain",
        "AI reasoning captured",
        "SOC 2 ready compliance",
    ], ACCENT_RED),
    ("Approval Workflows", [
        "Human-in-the-loop gates",
        "Auto-approve: high confidence + low risk",
        "Sensitive path detection",
        "TTL-based expiration",
    ], ACCENT_GREEN),
    ("Multi-Tenant Manager", [
        "Tenant-scoped data isolation",
        "Independent budget/SLA tracking",
        "Per-tenant: audit/, costs/, roi/",
        "Configurable allowed_repos",
    ], DARK_BLUE),
    ("Secrets Manager", [
        "Azure Key Vault + env fallback",
        "DefaultAzureCredential (MI)",
        "Rotation age monitoring (90d)",
        "Access audit (values never logged)",
    ], ACCENT_TEAL),
]

for i, (name, items, color) in enumerate(enterprise_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.7)
    add_card(slide, x, y, Inches(3.9), Inches(2.4), name, items, accent=color)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — INTELLIGENCE & FEEDBACK
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Intelligence & Feedback Loop",
               "5-factor business impact scoring + continuous learning from PR outcomes")

# Left: Business Impact Scoring
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.4),
             "Business Impact Scorer (5 Factors)", font_size=20, bold=True, color=DARK_BLUE)

factors = [
    ("Usage Frequency", "\u00d7100", "Telemetry-based call counts"),
    ("Criticality", "\u00d750", "Critical path: payments, auth, security"),
    ("Change Frequency", "\u00d730", "Git history churn rate"),
    ("Velocity Impact", "\u00d740", "Azure DevOps blocked work items"),
    ("Regulatory Risk", "\u00d780", "PII / GDPR / HIPAA annotations"),
]

for i, (name, weight, desc) in enumerate(factors):
    y = Inches(2.2 + i * 0.7)
    add_shape_bg(slide, Inches(0.9), y, Inches(1), Inches(0.5), AZURE_BLUE)
    add_text_box(slide, Inches(0.9), y + Inches(0.05), Inches(1), Inches(0.4),
                 weight, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), y + Inches(0.05), Inches(4), Inches(0.4),
                 f"{name} — {desc}", font_size=14, color=DARK_GRAY)

# Right: Feedback & Learning
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
             "Feedback & Learning System", font_size=20, bold=True, color=DARK_BLUE)

feedback_items = [
    "PR Outcome Tracking: merged / rejected / modified",
    "Per-scanner success rates with rolling averages",
    "Auto-adjusting confidence thresholds",
    "Team preference learning \u2192 injected into system prompt",
    "Historical pattern recognition across org",
    "Similarity scoring for proven refactoring approaches",
]
add_bullet_list(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.5),
                ["  " + f for f in feedback_items], font_size=14, color=DARK_GRAY, bold_prefix=False)

# Dynamic reprioritizer
add_text_box(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(0.4),
             "Dynamic Reprioritizer", font_size=18, bold=True, color=DARK_BLUE)

events = [
    "Production incidents \u2192 boost affected findings",
    "New CVE announcements \u2192 boost security findings",
    "Sprint deadline proximity \u2192 boost velocity impact",
    "Budget changes \u2192 reorder by cost-effectiveness",
]
add_bullet_list(slide, Inches(7), Inches(5.1), Inches(5.5), Inches(2),
                ["  " + e for e in events], font_size=14, color=DARK_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CI SELF-HEALING & PR REVIEW BOT
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "CI Self-Healing & PR Review Bot",
               "Unique capabilities — no competitor has these")

# Left: CI Self-Healing
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(5.9), Inches(5.0), LIGHT_GRAY)
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(0.08), Inches(5.0), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4),
             "CI Self-Healing (First in Market)", font_size=20, bold=True, color=ACCENT_RED)

heal_steps = [
    "GitHub Action triggers on CI failure",
    "Reads failure log, extracts typed failure signals",
    "Pattern-matches: Ruff F401, mypy incompatible-return, pytest assertion",
    "Generates patch candidates for each failure signal",
    "Copilot SDK generates complex fixes beyond pattern matching",
    "Posts idempotent healing plan as PR comment",
    "Prevents duplicate comments via HTML marker",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.3), Inches(3.5),
                ["  " + s for s in heal_steps], font_size=14, color=DARK_GRAY, bold_prefix=False)

# Right: PR Review Bot
add_shape_bg(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.0), LIGHT_GRAY)
add_shape_bg(slide, Inches(6.9), Inches(1.6), Inches(0.08), Inches(5.0), AZURE_BLUE)

add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.4),
             "PR Review Bot", font_size=20, bold=True, color=AZURE_BLUE)

review_items = [
    "Scans incoming PRs for code smells, security, deprecated APIs",
    "Posts structured risk summary as PR comment",
    "Auto-labels: needs-fix, security-risk, type-issues",
    "Configurable --block-on severity gating",
    "Copilot SDK generates fix suggestions in review",
    "Triggers on pull_request + workflow_dispatch",
    "CodeCustodian reviews AND can fix — CodeRabbit only reviews",
]
add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5.3), Inches(3.5),
                ["  " + r for r in review_items], font_size=14, color=DARK_GRAY, bold_prefix=False)

add_text_box(slide, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.5),
             "Dependabot creates PRs that break builds and walks away. CodeCustodian fixes its own CI failures — true self-healing.",
             font_size=15, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — DEPLOYMENT ARCHITECTURE (AZURE)
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Deployment Architecture — Azure",
               "Infrastructure as Code: main.bicep + 8 modules")

# Azure resources
resources = [
    ("Managed Identity", "Zero-secret\nauthentication", AZURE_BLUE),
    ("Container Registry", "Image storage\nwith RBAC", ACCENT_TEAL),
    ("Key Vault", "Secret management\n+ rotation", DARK_BLUE),
    ("Container Apps", "Runtime hosting\nauto-scaling", ACCENT_GREEN),
    ("VNet + Subnets", "Network isolation\nfor workloads", MED_GRAY),
    ("Log Analytics", "Centralized\nlog aggregation", ACCENT_ORANGE),
    ("App Insights", "APM + distributed\ntracing", AZURE_BLUE),
    ("Dashboard + Alerts", "Operational visibility\n+ proactive monitoring", ACCENT_RED),
]

for i, (name, desc, color) in enumerate(resources):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.6 + row * 2.1)
    box = add_shape_bg(slide, x, y, Inches(2.85), Inches(1.8), color)
    add_text_box(slide, x, y + Inches(0.15), Inches(2.85), Inches(0.4),
                 name, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.6), Inches(2.85), Inches(0.9),
                 desc, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

# Deployment flow
add_text_box(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.4),
             "Deployment Flow:  Docker Build  \u2192  ACR Push  \u2192  Bicep Deploy  \u2192  Container App Revision  \u2192  Health Check  \u2192  Traffic Shift",
             font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
             "Transport: Streamable HTTP on port 8080  |  Secrets injected from Key Vault via Managed Identity  |  Auto-scaling based on HTTP traffic",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — SECURITY ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Security Architecture",
               "Defense-in-depth: 5 layers from authentication to infrastructure")

layers = [
    ("Layer 1: Authentication", "Azure AD / JWT tokens + Managed Identity\nfor all inter-service communication", AZURE_BLUE),
    ("Layer 2: Authorization", "RBAC engine (6 roles x 10 permissions)\n+ multi-tenant data isolation", ACCENT_TEAL),
    ("Layer 3: Code Safety", "6 pre-execution safety checks\n+ 9 regex secret detection patterns", ACCENT_ORANGE),
    ("Layer 4: Data Integrity", "Audit logger with SHA-256 tamper-evident\nhashes + atomic backup/rollback", ACCENT_GREEN),
    ("Layer 5: Infrastructure", "Key Vault secret rotation + VNet isolation\n+ TLS termination at ingress", DARK_BLUE),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(1.6 + i * 1.1)
    box = add_shape_bg(slide, Inches(0.5), y, Inches(12.3), Inches(0.9), color)
    add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(4), Inches(0.7),
                 title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(4.8), y + Inches(0.1), Inches(7.8), Inches(0.7),
                 desc, font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — OBSERVABILITY
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Observability Architecture",
               "Full OpenTelemetry integration with Azure Monitor")

# Left: Telemetry types
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.4),
             "Telemetry Types", font_size=20, bold=True, color=DARK_BLUE)

tel = [
    "Traces: Pipeline spans, scanner spans, planner spans — full distributed tracing",
    "Metrics: findings_count, pr_created, duration_seconds, cost_usd — custom counters",
    "Logs: Structured JSON via get_logger() with contextual fields",
    "Export: Azure Monitor (production) + Console (development) — dual exporter",
    "Graceful fallback when Azure SDK is not installed — local logging only",
]
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.5), Inches(3),
                ["  " + t for t in tel], font_size=14, color=DARK_GRAY, bold_prefix=False)

# Right: Span hierarchy
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
             "Span Hierarchy", font_size=20, bold=True, color=DARK_BLUE)

spans = [
    "pipeline.run (root)",
    "  \u251C\u2500 pipeline.scan",
    "  \u2502   \u251C\u2500 pipeline.scan.deprecated_api",
    "  \u2502   \u251C\u2500 pipeline.scan.todo_comments",
    "  \u2502   \u251C\u2500 pipeline.scan.code_smells",
    "  \u2502   \u2514\u2500 pipeline.scan.security",
    "  \u251C\u2500 pipeline.dedup",
    "  \u251C\u2500 pipeline.prioritize",
    "  \u251C\u2500 pipeline.plan.{finding_id}",
    "  \u251C\u2500 pipeline.execute.{finding_id}",
    "  \u251C\u2500 pipeline.verify.{finding_id}",
    "  \u2514\u2500 pipeline.pr.{finding_id}",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4),
                spans, font_size=12, color=DARK_GRAY, bold_prefix=False,
                spacing=Pt(2))

# SLA Reporting
add_text_box(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.4),
             "SLA Reporting", font_size=18, bold=True, color=DARK_BLUE)
sla = [
    "Success rate, avg duration, P95 duration, avg time-to-PR",
    "Top failure reasons with trend analysis (improving / stable / degrading)",
    "Failure spike detection with configurable threshold (default 10%)",
]
add_bullet_list(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.5),
                ["  " + s for s in sla], font_size=14, color=DARK_GRAY, bold_prefix=False)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — COMPETITIVE ADVANTAGES
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Competitive Advantages",
               "What no other tool delivers")

comparisons = [
    ("vs SonarQube", "SonarQube detects issues. CodeCustodian detects,\nplans with AI, executes, verifies, and creates PRs.", AZURE_BLUE),
    ("vs Dependabot", "Dependabot bumps version numbers.\nCodeCustodian bumps AND migrates code in one PR.", ACCENT_TEAL),
    ("vs CodeRabbit", "CodeRabbit reviews PRs (reactive).\nCodeCustodian creates PRs proactively + reviews.", ACCENT_GREEN),
    ("vs Moderne", "Moderne requires hand-written Java recipes.\nCodeCustodian uses AI — no recipe authoring needed.", DARK_BLUE),
    ("vs Snyk", "Snyk is security-only SaaS ($$$).\nCodeCustodian: open-source, 6 scanner types, AI fixes.", ACCENT_ORANGE),
    ("vs Everyone", "No competitor has CI self-healing, MCP server,\nWork IQ org context, or feedback-driven learning.", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(comparisons):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.35)
    y = Inches(1.6 + row * 1.8)
    add_card(slide, x, y, Inches(6.1), Inches(1.55), title, desc.split("\n"), accent=color)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — SDK INTEGRATION MAP
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "SDK & Technology Stack",
               "13 SDK integrations powering the autonomous pipeline")

sdk_groups = [
    ("AI & MCP", [
        "GitHub Copilot SDK (copilot >=0.1.29): Sessions, tools, model routing, cost tracking",
        "FastMCP v2 (fastmcp >=2.14.0): Server (8 tools, 7 resources, 4 prompts) + Client (Work IQ)",
    ], AZURE_BLUE),
    ("Azure SDKs", [
        "azure-identity: DefaultAzureCredential for Key Vault + Managed Identity",
        "azure-keyvault-secrets: SecretClient for secret retrieval + rotation monitoring",
        "azure-monitor-opentelemetry: One-liner configure_azure_monitor() setup",
        "httpx.AsyncClient: Azure DevOps REST API for work items",
    ], ACCENT_TEAL),
    ("GitHub & Git", [
        "PyGithub: PR creation, inline comments, issue creation, label management",
        "GitPython: Branch ops, conventional commits, push, SHA tracking",
    ], DARK_BLUE),
    ("Data & Runtime", [
        "Pydantic v2 (>=2.5): All models + config schema + validation",
        "OpenTelemetry: Distributed tracing spans across pipeline",
        "TinyDB: Lightweight JSON DB for SLA, budget, feedback",
        "Typer + Rich: CLI framework with beautiful console output",
    ], ACCENT_GREEN),
]

y = Inches(1.6)
for title, items, color in sdk_groups:
    h = Inches(0.5 + len(items) * 0.45)
    add_card(slide, Inches(0.5), y, Inches(12.3), h, title, items, accent=color)
    y += h + Inches(0.15)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Configuration Architecture",
               ".codecustodian.yml — Pydantic v2 validated, 14 config sections")

# Left: hierarchy
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.4),
             "Configuration Hierarchy", font_size=20, bold=True, color=DARK_BLUE)

hierarchy = [
    "Organization defaults (global baseline)",
    "  \u2192 Team overrides (per-team customization)",
    "    \u2192 Repository .codecustodian.yml (repo-specific)",
    "      \u2192 CLI flags (runtime override, highest priority)",
]
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.5), Inches(1.5),
                hierarchy, font_size=15, color=DARK_GRAY, bold_prefix=False)

# Config sections
sections = [
    "scanners \u2014 Per-scanner config (deprecated_apis, todo, smells, security, types, deps)",
    "behavior \u2014 max_prs_per_run, confidence_threshold, proposal_mode_threshold",
    "github \u2014 repo_name, pr_labels, reviewers, base_branch, draft_threshold",
    "copilot \u2014 model_strategy (auto/fast/balanced/reasoning), BYOK, alternatives",
    "budget \u2014 monthly_budget, alerts at threshold percentages",
    "approval \u2014 require_approval, auto_approve criteria, sensitive paths",
    "sla \u2014 enabled, db_path, failure_spike_threshold",
    "work_iq \u2014 enabled, command, args for MCP server",
]
add_bullet_list(slide, Inches(0.9), Inches(3.7), Inches(11.5), Inches(3.5),
                ["  " + s for s in sections], font_size=14, color=DARK_GRAY, bold_prefix=True)

# Validation callout
add_text_box(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
             "Validation: @field_validator for range/format  |  @model_validator for cross-field consistency  |  ConfigDict(validate_assignment=True) for runtime safety",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19 — BY THE NUMBERS
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT_TEAL)

add_text_box(slide, Inches(0.9), Inches(0.4), Inches(11), Inches(0.7),
             "CodeCustodian by the Numbers", font_size=36, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

stats = [
    ("7", "Scanners"),
    ("13", "AI Models"),
    ("8", "MCP Tools"),
    ("7", "MCP Resources"),
    ("6", "Safety Checks"),
    ("5", "Impact Factors"),
    ("6", "RBAC Roles"),
    ("8", "Bicep Modules"),
    ("14", "Config Sections"),
    ("80%+", "Test Coverage"),
    ("95%", "PR Acceptance"),
    ("20+", "hrs/wk Saved"),
]

for i, (num, label) in enumerate(stats):
    col = i % 6
    row = i // 6
    x = Inches(0.5 + col * 2.1)
    y = Inches(1.5 + row * 2.8)
    box = add_shape_bg(slide, x, y, Inches(1.85), Inches(2.2), AZURE_BLUE)
    add_text_box(slide, x, y + Inches(0.3), Inches(1.85), Inches(0.8),
                 num, font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.2), Inches(1.85), Inches(0.6),
                 label, font_size=15, color=RGBColor(0xBB, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20 — THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT_TEAL)
add_shape_bg(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.2),
             "Thank You", font_size=60, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

add_text_box(slide, Inches(2), Inches(3.5), Inches(9.5), Inches(0.7),
             "Questions & Discussion",
             font_size=28, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

links = [
    "GitHub: github.com/Nizarel/CodeCustodian",
    "Powered by: GitHub Copilot SDK  |  FastMCP v2  |  Azure Container Apps",
    "Python 3.11+  |  Pydantic v2  |  OpenTelemetry  |  Bicep IaC",
]
add_bullet_list(slide, Inches(2.5), Inches(4.6), Inches(8.5), Inches(2),
                links, font_size=16, color=RGBColor(0xAA, 0xCC, 0xEE),
                bold_prefix=False, spacing=Pt(12))


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════

output_path = r"c:\Users\nizare\OneDrive - Microsoft\Documents\GitHub\CodeCustodian\presentations\CodeCustodian-Features.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
